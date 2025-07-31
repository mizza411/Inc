"""
Automated Pitch Deck Generator
Creates pitch decks following the proven template structure with business validation data.
"""

import os
import subprocess
import platform
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PitchDeckGenerator:
    def __init__(self, template_path="Pitch Deck (YOURS).. (2).pptx"):
        """Initialize with the template path."""
        self.template_path = template_path
        self.output_dir = "generated_pitch_decks"
        
        # Create output directory if it doesn't exist
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)
    
    def generate_pitch_deck(self, gadget_data):
        """Generate a pitch deck for a specific gadget using the template."""
        try:
            # Load the template
            prs = Presentation(self.template_path)
            
            # Fill in the slides with gadget data
            self._fill_title_slide(prs.slides[0], gadget_data)
            self._fill_credibility_slide(prs.slides[1], gadget_data)
            self._fill_problem_slide(prs.slides[2], gadget_data)
            self._fill_market_potential_slide(prs.slides[3], gadget_data)
            self._fill_current_solutions_slide(prs.slides[4], gadget_data)
            self._fill_proposed_solution_slide(prs.slides[5], gadget_data)
            self._fill_dream_slide(prs.slides[6], gadget_data)
            self._fill_traction_slide(prs.slides[7], gadget_data)
            self._fill_solution_details_slide(prs.slides[8], gadget_data)
            self._fill_differentiation_slide(prs.slides[9], gadget_data)
            self._fill_financials_slide(prs.slides[10], gadget_data)
            self._fill_trend_slide(prs.slides[11], gadget_data)
            self._fill_why_now_slide(prs.slides[12], gadget_data)
            self._fill_why_you_slide(prs.slides[13], gadget_data)
            self._fill_milestones_slide(prs.slides[14], gadget_data)
            self._fill_pot_sweetener_slide(prs.slides[15], gadget_data)
            self._fill_call_to_action_slide(prs.slides[16], gadget_data)
            
            # Generate filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"pitch_deck_{gadget_data['name'].replace(' ', '_')}_{timestamp}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            
            # Save the presentation
            prs.save(filepath)
            
            return filepath
            
        except Exception as e:
            print(f"Error generating pitch deck: {e}")
            return None
    
    def _fill_title_slide(self, slide, data):
        """Fill the title slide with gadget information."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = f"{data['name']} Business"
        
        # Add subtitle with validation status
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = f"Validated Business Opportunity\n{data.get('validation_status', 'Validated')}"
    
    def _fill_credibility_slide(self, slide, data):
        """Fill credibility slide with validation data."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Business Validation Credibility"
        
        # Add validation metrics
        content = f"""
        ✅ DEMAND VALIDATION: PASSED
        • Demand Score: {data.get('demand_score', 0)}/100
        • Market Trend: {data.get('trend', 'Increasing')}
        
        ✅ MARKET VALIDATION: PASSED
        • Available Suppliers: {len(data.get('sourcing_options', []))}
        • Best Source: {data.get('best_source', 'Alibaba')}
        
        ✅ COMPETITIVE ADVANTAGE: CONFIRMED
        • Viability Score: {data.get('viability_score', 0)}/100
        • Status: {data.get('validation_status', 'Validated')}
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_problem_slide(self, slide, data):
        """Fill problem slide with market opportunity."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Market Problem & Opportunity"
        
        content = f"""
        PROBLEM:
        • High demand for {data['name']} with limited local supply
        • Price variations indicate market inefficiency
        • Consumers seeking quality at competitive prices
        
        OPPORTUNITY:
        • {data.get('demand_score', 0)}/100 demand score
        • {data.get('trend', 'Increasing')} market trend
        • {len(data.get('sourcing_options', []))} supplier options available
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_market_potential_slide(self, slide, data):
        """Fill market potential slide with demand analysis."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Market Potential Analysis"
        
        content = f"""
        MARKET SIZE & DEMAND:
        • High demand score: {data.get('demand_score', 0)}/100
        • Market trend: {data.get('trend', 'Increasing')}
        • Multiple supplier options: {len(data.get('sourcing_options', []))}
        
        COMPETITIVE LANDSCAPE:
        """
        
        # Safely calculate price range
        sourcing_options = data.get('sourcing_options', [])
        prices = []
        for option in sourcing_options:
            price_str = option.get('price', 'N/A')
            if price_str != 'N/A':
                try:
                    import re
                    clean_price = re.sub(r'[₦$,\s]', '', price_str)
                    if clean_price and clean_price.replace('.', '').isdigit():
                        prices.append(float(clean_price))
                except (ValueError, AttributeError):
                    continue
        
        if prices:
            content += f"• Price range: ${min(prices):.2f} - ${max(prices):.2f}\n"
            content += "• Significant price variation indicates market opportunity"
        else:
            content += "• Multiple pricing options available\n"
            content += "• Market opportunity confirmed through supplier diversity"
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_current_solutions_slide(self, slide, data):
        """Fill current solutions slide with competitive analysis."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Current Market Solutions"
        
        sourcing_options = data.get('sourcing_options', [])
        content = "CURRENT SUPPLIERS:\n\n"
        
        for i, option in enumerate(sourcing_options[:3], 1):  # Show top 3
            content += f"{i}. {option.get('platform', 'Unknown')}\n"
            content += f"   Price: {option.get('price', 'N/A')}\n"
            content += f"   Rating: {option.get('rating', 'N/A')}\n\n"
        
        # Calculate price variation safely
        prices = []
        for option in sourcing_options:
            price_str = option.get('price', 'N/A')
            if price_str != 'N/A':
                try:
                    # Remove currency symbols and convert to float
                    import re
                    # Remove ₦, $, and other currency symbols
                    clean_price = re.sub(r'[₦$,\s]', '', price_str)
                    if clean_price and clean_price.replace('.', '').isdigit():
                        prices.append(float(clean_price))
                except (ValueError, AttributeError):
                    continue
        
        if prices:
            price_range = max(prices) - min(prices)
            content += f"OPPORTUNITY: Price variation of ${price_range:.2f} indicates market inefficiency"
        else:
            content += "OPPORTUNITY: Multiple supplier options available with competitive pricing"
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_proposed_solution_slide(self, slide, data):
        """Fill proposed solution slide with business model."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Our Proposed Solution"
        
        content = f"""
        BUSINESS MODEL:
        • Source from {data.get('best_source', 'Alibaba')} (B2B pricing)
        • Sell at competitive retail prices
        • Target {data.get('profit_margin', 60):.1f}% profit margin
        
        VALIDATED APPROACH:
        • {data.get('validation_status', 'Validated')} business model
        • {data.get('viability_score', 0)}/100 viability score
        • {data.get('demand_score', 0)}/100 demand score
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_dream_slide(self, slide, data):
        """Fill dream slide with vision."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Our Vision"
        
        content = f"""
        THE DREAM:
        • Become the leading supplier of {data['name']} in the market
        • Leverage validated demand of {data.get('demand_score', 0)}/100
        • Capture market share through competitive pricing
        • Build sustainable business with {data.get('profit_margin', 60):.1f}% margins
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_traction_slide(self, slide, data):
        """Fill traction slide with validation evidence."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Evidence of Traction"
        
        content = f"""
        VALIDATION EVIDENCE:
        ✅ Demand Score: {data.get('demand_score', 0)}/100
        ✅ Market Trend: {data.get('trend', 'Increasing')}
        ✅ Supplier Availability: {len(data.get('sourcing_options', []))} sources
        ✅ Competitive Advantage: Confirmed
        ✅ Viability Score: {data.get('viability_score', 0)}/100
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_solution_details_slide(self, slide, data):
        """Fill solution details slide with implementation plan."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Solution Implementation Details"
        
        content = f"""
        IMPLEMENTATION PLAN:
        1. Source from {data.get('best_source', 'Alibaba')} at B2B prices
        2. Establish competitive retail pricing strategy
        3. Target {data.get('profit_margin', 60):.1f}% profit margins
        4. Scale based on {data.get('demand_score', 0)}/100 demand score
        
        VALIDATED APPROACH:
        • {data.get('validation_status', 'Validated')} business model
        • {data.get('viability_score', 0)}/100 viability score
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_differentiation_slide(self, slide, data):
        """Fill differentiation slide with competitive advantages."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "What Differentiates Our Solution"
        
        content = f"""
        COMPETITIVE ADVANTAGES:
        • B2B sourcing from {data.get('best_source', 'Alibaba')}
        • {data.get('profit_margin', 60):.1f}% profit margins
        • Validated demand of {data.get('demand_score', 0)}/100
        • {data.get('viability_score', 0)}/100 viability score
        
        MARKET POSITIONING:
        • Competitive pricing through B2B relationships
        • Quality assurance through validated suppliers
        • Scalable model based on market demand
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_financials_slide(self, slide, data):
        """Fill financials slide with business model and projections."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Financials and Business Model"
        
        # Calculate financial projections
        product_cost = data.get('product_cost', 15.50)
        retail_price = data.get('retail_price', 38.75)
        profit_margin = data.get('profit_margin', 60)
        funding_needed = data.get('funding_needed', 4456.25)
        
        content = f"""
        BUSINESS MODEL:
        • Product Cost: ${product_cost:.2f}
        • Retail Price: ${retail_price:.2f}
        • Profit Margin: {profit_margin:.1f}%
        
        FUNDING REQUIREMENTS:
        • Total Funding Needed: ${funding_needed:.2f}
        • Initial Inventory: ${product_cost * 100:.2f}
        • Marketing Budget: ${retail_price * 50:.2f}
        • Operational Costs: ${retail_price * 25:.2f}
        
        PROJECTED RETURNS:
        • ROI based on {profit_margin:.1f}% margins
        • Scalable model with validated demand
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_trend_slide(self, slide, data):
        """Fill trend slide with market trends."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Larger Market Trend"
        
        content = f"""
        MARKET TRENDS:
        • {data.get('trend', 'Increasing')} demand for {data['name']}
        • {data.get('demand_score', 0)}/100 demand score
        • Multiple supplier options available
        • Price variations indicate market opportunity
        
        INDUSTRY INSIGHTS:
        • Growing market for smart gadgets
        • B2B sourcing becoming more accessible
        • Consumer demand for competitive pricing
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_why_now_slide(self, slide, data):
        """Fill why now slide with timing factors."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Why Now?"
        
        content = f"""
        TIMING FACTORS:
        • High demand score: {data.get('demand_score', 0)}/100
        • {data.get('trend', 'Increasing')} market trend
        • {data.get('validation_status', 'Validated')} business model
        • {data.get('viability_score', 0)}/100 viability score
        
        MARKET CONDITIONS:
        • Multiple supplier options available
        • Price variations create opportunity
        • B2B sourcing accessible
        • Validated demand exists
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_why_you_slide(self, slide, data):
        """Fill why you slide with team capabilities."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Why You?"
        
        content = f"""
        TEAM CAPABILITIES:
        • Data-driven approach with {data.get('demand_score', 0)}/100 demand validation
        • B2B sourcing expertise from {data.get('best_source', 'Alibaba')}
        • {data.get('profit_margin', 60):.1f}% profit margin strategy
        • {data.get('viability_score', 0)}/100 viability score
        
        COMPETITIVE ADVANTAGES:
        • Validated business model
        • Objective data-driven decisions
        • Proven template-based approach
        • Scalable operations
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_milestones_slide(self, slide, data):
        """Fill milestones slide with future goals."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Milestones and Future Goals"
        
        content = f"""
        SHORT-TERM GOALS (3-6 months):
        • Secure ${data.get('funding_needed', 4456.25):.2f} in funding
        • Establish supplier relationships
        • Launch initial inventory
        
        MEDIUM-TERM GOALS (6-12 months):
        • Scale based on {data.get('demand_score', 0)}/100 demand
        • Achieve {data.get('profit_margin', 60):.1f}% profit margins
        • Expand product line
        
        LONG-TERM VISION:
        • Market leadership in {data['name']} category
        • Sustainable {data.get('profit_margin', 60):.1f}% profit model
        • Validated business expansion
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_pot_sweetener_slide(self, slide, data):
        """Fill pot sweetener slide with additional benefits."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Additional Benefits"
        
        content = f"""
        POT SWEETENER:
        • {data.get('validation_status', 'Validated')} business model
        • {data.get('viability_score', 0)}/100 viability score
        • {data.get('demand_score', 0)}/100 demand validation
        • {data.get('profit_margin', 60):.1f}% profit margins
        
        RISK MITIGATION:
        • Objective data validation
        • Multiple supplier options
        • Proven template approach
        • Scalable business model
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content
    
    def _fill_call_to_action_slide(self, slide, data):
        """Fill call to action slide with next steps."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Call to Action"
        
        content = f"""
        NEXT STEPS:
        • Invest ${data.get('funding_needed', 4456.25):.2f} in validated opportunity
        • Support {data.get('validation_status', 'Validated')} business model
        • Join {data.get('viability_score', 0)}/100 viability venture
        
        CONTACT:
        • Generated pitch deck: {data.get('pitch_deck_path', 'Available')}
        • Validation data: {data.get('validation_status', 'Available')}
        • Business model: {data.get('profit_margin', 60):.1f}% margins
        """
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                shape.text_frame.text = content

def open_file_automatically(filepath):
    """Automatically open a file using the default system application."""
    try:
        if platform.system() == "Windows":
            os.startfile(filepath)
        elif platform.system() == "Darwin":  # macOS
            subprocess.run(["open", filepath])
        else:  # Linux
            subprocess.run(["xdg-open", filepath])
        return True
    except Exception as e:
        print(f"Could not open file automatically: {e}")
        return False

def generate_pitch_deck_for_gadget(gadget_data):
    """Generate a pitch deck for a specific gadget."""
    generator = PitchDeckGenerator()
    filepath = generator.generate_pitch_deck(gadget_data)
    
    if filepath:
        print(f"📄 Pitch deck generated: {filepath}")
        print(f"🚀 Opening pitch deck automatically...")
        
        # Automatically open the file
        if open_file_automatically(filepath):
            print(f"✅ Pitch deck opened successfully!")
        else:
            print(f"⚠️  Could not open automatically. Please open manually: {filepath}")
        
        return filepath
    else:
        print("❌ Failed to generate pitch deck")
        return None 